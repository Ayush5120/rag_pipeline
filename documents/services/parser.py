import io 

def extract_text(file_field) -> str:
    """
    Accepts a Django FieldFile (doc.file) and returns plain text.
    Supports: .txt, .pdf, .docx, .pptx, .xlsx/.csv
    """
    filename = file_field.name.lower()
    file_field.seek(0)
    raw = file_field.read()

    if filename.endswith('.txt') or filename.endswith('.md'):
        return _read_txt(raw)

    elif filename.endswith('.pdf'):
        return _read_pdf(raw)

    elif filename.endswith('.docx'):
        return _read_docx(raw)

    elif filename.endswith('.pptx'):
        return _read_pptx(raw)

    elif filename.endswith('.xlsx'):
        return _read_xlsx(raw)

    elif filename.endswith('.csv'):
        return _read_csv(raw)

    else:
        raise ValueError(f"Unsupported file type: {file_field.name}")
    

def _read_txt(raw: bytes) -> str:
    return raw.decode('utf-8', errors='replace')

def _read_pdf(raw: bytes) -> str:
    import pdfplumber  # pip install pdfplumber

    text_parts = []
    with pdfplumber.open(io.BytesIO(raw)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)

    result =  '\n\n'.join(text_parts)

    if not result.strip():
        raise ValueError(
            "No text extracted from PDF. "
            "This may be a scanned/image-only PDF requiring OCR. "
            "Try uploading a text-based PDF."
        )
    return result

def _read_docx(raw: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(raw))
    parts = []

    # Body paragraphs
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)

    # Tables — convert to readable text
    for table in doc.tables:
        for row in table.rows:
            row_text = ' | '.join(
                cell.text.strip()
                for cell in row.cells
                if cell.text.strip()
            )
            if row_text:
                parts.append(row_text)

    return '\n\n'.join(parts)


def _read_pptx(raw: bytes) -> str:
    from pptx import Presentation  # pip install python-pptx

    prs = Presentation(io.BytesIO(raw))
    slide_texts = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            slide_texts.append(f"[Slide {i+1}]\n" + '\n'.join(texts))
    return '\n\n'.join(slide_texts)
    # [Slide N] prefix helps the chunker keep context of where text came from.


def _read_xlsx(raw: bytes) -> str:
    import openpyxl  # pip install openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    sheet_texts = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            row_text = '\t'.join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                rows.append(row_text)
        if rows:
            sheet_texts.append(f"[Sheet: {sheet.title}]\n" + '\n'.join(rows))
    return '\n\n'.join(sheet_texts)
    # Tab-separated rows = readable for the LLM.
    # data_only=True returns computed values, not formulas.


def _read_csv(raw: bytes) -> str:
    import csv

    # Try UTF-8 first, fall back to latin-1
    # Excel on Windows saves CSVs as latin-1 by default
    try:
        text = raw.decode('utf-8')
    except UnicodeDecodeError:
        text = raw.decode('latin-1')
    reader = csv.reader(io.StringIO(text))
    rows = ['\t'.join(row) for row in reader if any(cell.strip() for cell in row)]
    return '\n'.join(rows)
